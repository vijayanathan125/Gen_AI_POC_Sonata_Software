import os
import io
import PyPDF2
import openai
import openpyxl
from docx import Document
from pptx import Presentation
import ezodf
import pandas as pd
import torch
from transformers import BertTokenizer, BertConfig, BertForSequenceClassification, AdamW
from torch.utils.data import DataLoader, Dataset
from sklearn.model_selection import train_test_split
from azure.storage.blob import BlobServiceClient
from flask import Flask, jsonify, request
from flask_cors import CORS

app = Flask(__name__)
CORS(app, resources={"/bert_query": {"origins": "http://localhost:3000"}})

# Set your Azure Storage account connection string
azure_storage_connection_string = "AccountName=devstoreaccount1;AccountKey=Eby8vdM02xNOcqFlqUwJPLlmEtlCDXJ1OUzFT50uSRZ6IFsuFq2UVErCz4I6tq/K1SZFPTOtr/KBHBeksoGMGw==;DefaultEndpointsProtocol=http;BlobEndpoint=http://127.0.0.1:10000/devstoreaccount1;QueueEndpoint=http://127.0.0.1:10001/devstoreaccount1;TableEndpoint=http://127.0.0.1:10002/devstoreaccount1;"

# Initialize BlobServiceClient
blob_service_client = BlobServiceClient.from_connection_string(azure_storage_connection_string)

# Load and preprocess your dataset from Azure Blob Storage
container_client = blob_service_client.get_container_client("uploaded-documents")
blob_client = container_client.get_blob_client("training_set.csv")

# Download dataset file from Azure Blob Storage
downloaded_blob = blob_client.download_blob()
df = pd.read_csv(io.BytesIO(downloaded_blob.readall()), encoding='latin1')
df = df.dropna()  # Drop rows with missing values
df = df.sample(frac=1).reset_index(drop=True)  # Shuffle the dataset

# Split the dataset into training and validation sets
df = df.sample(frac=1, random_state=42).reset_index(drop=True)
train_df, val_df = train_test_split(df, test_size=0.2, random_state=42)

# Define a custom dataset class
class CustomDataset(Dataset):
    def __init__(self, texts, labels, tokenizer, max_len=128):
        self.texts = texts
        self.labels = labels
        self.tokenizer = tokenizer
        self.max_len = max_len

    def __len__(self):
        return len(self.texts)

    def __getitem__(self, idx):
        text = str(self.texts.iloc[idx])
        label = int(self.labels.iloc[idx])
        inputs = self.tokenizer(text, padding=True, truncation=True, max_length=self.max_len, return_tensors='pt')

        return {
            'text': text,
            'label': label,
            'input_ids': inputs['input_ids'].squeeze(),
            'attention_mask': inputs['attention_mask'].squeeze(),
        }

# Set your OpenAI GPT-3.5 API key
api_key = "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"
openai.api_key = api_key

# Initialize the BERT tokenizer and model
tokenizer = BertTokenizer.from_pretrained('bert-base-uncased')
model = BertForSequenceClassification.from_pretrained('bert-base-uncased', num_labels=2)  # Change num_labels according to your task
model.eval()

# Prepare the training and validation datasets
train_dataset = CustomDataset(train_df['text'], train_df['label'].astype(int), tokenizer)
val_dataset = CustomDataset(val_df['text'], val_df['label'].astype(int), tokenizer)

# DataLoader for efficient loading of data during training
def collate_fn(batch):
    texts = [item['text'] for item in batch]
    labels = [item['label'] for item in batch]

    # Tokenize and pad the sequences
    inputs = tokenizer(texts, padding=True, truncation=True, return_tensors="pt")
    input_ids = inputs["input_ids"]
    attention_mask = inputs["attention_mask"]

    return {
        "input_ids": input_ids,
        "attention_mask": attention_mask,
        "labels": torch.tensor(labels, dtype=torch.long)
    }

train_dataloader = DataLoader(train_dataset, batch_size=8, shuffle=True, collate_fn=collate_fn)
val_dataloader = DataLoader(val_dataset, batch_size=8, shuffle=False, collate_fn=collate_fn)

# Define optimizer and training parameters
optimizer = AdamW(model.parameters(), lr=2e-5)
num_epochs = 3

# Define the function to save the model locally
def save_model_locally(model, config, output_dir):
    # Create the directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)
    # Save model weights
    model.save_pretrained(output_dir)
    # Save model configuration
    config.save_pretrained(output_dir)

# Training loop
device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
model.to(device)

for epoch in range(num_epochs):
    model.train()
    train_loss = 0.0
    for batch in train_dataloader:
        inputs = {k: v.to(device) for k, v in batch.items()}
        optimizer.zero_grad()
        outputs = model(**inputs, return_dict=True)
        loss = outputs.loss
        loss.backward()
        optimizer.step()
        train_loss += loss.item()

    avg_train_loss = train_loss / len(train_dataloader)

    # Validation loop
    model.eval()
    val_loss = 0.0
    with torch.no_grad():
        for batch in val_dataloader:
            inputs = {k: v.to(device) for k, v in batch.items()}
            outputs = model(**inputs, return_dict=True)
            loss = outputs.loss
            val_loss += loss.item()

    avg_val_loss = val_loss / len(val_dataloader)

    print(f'Epoch {epoch + 1}/{num_epochs}, Training Loss: {avg_train_loss:.4f}, Validation Loss: {avg_val_loss:.4f}')

# Save the fine-tuned BERT model locally
model_output_dir = "fine_tuned_document_search_model"
save_model_locally(model, tokenizer, model_output_dir)

# Load the fine-tuned BERT model from the local file system
fine_tuned_model_path = model_output_dir
config = BertConfig.from_pretrained(fine_tuned_model_path)
fine_tuned_model = BertForSequenceClassification.from_pretrained(fine_tuned_model_path, config=config)

# Define other functions and routes as before...
# Function to read text content from different types of documents
# Function to read text content from different types of documents from Azure Blob Storage
def read_text_document(document_path):
    try:
        blob_client = container_client.get_blob_client(document_path)
        with io.BytesIO(blob_client.download_blob().readall()) as file:
            return file.read().decode('utf-8')
    except Exception as ex:
        return f"Error reading document: {str(ex)}"

# Function to read text content from a Word document from Azure Blob Storage
def read_word_document(document_path):
    try:
        blob_client = container_client.get_blob_client(document_path)
        with io.BytesIO(blob_client.download_blob().readall()) as file:
            doc = Document(io.BytesIO(file.read()))
            full_text = []
            for paragraph in doc.paragraphs:
                full_text.append(paragraph.text)
            return '\n'.join(full_text)
    except Exception as ex:
        return f"Error reading Word document: {str(ex)}"

# Function to read text content from a PDF document from Azure Blob Storage
def read_pdf_document(document_path):
    try:
        blob_client = container_client.get_blob_client(document_path)
        with io.BytesIO(blob_client.download_blob().readall()) as file:
            pdf_reader = PyPDF2.PdfReader(file)
            full_text = []
            for page_num, page in enumerate(pdf_reader.pages):
                full_text.append(page.extract_text())
            return '\n'.join(full_text)
    except Exception as ex:
        return f"Error reading PDF document: {str(ex)}"

# Function to read text content from a PowerPoint presentation from Azure Blob Storage
def read_powerpoint_document(document_path):
    try:
        blob_client = container_client.get_blob_client(document_path)
        with io.BytesIO(blob_client.download_blob().readall()) as file:
            presentation = Presentation(io.BytesIO(file.read()))
            full_text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        full_text.append(shape.text)
            return '\n'.join(full_text)
    except Exception as ex:
        return f"Error reading PowerPoint document: {str(ex)}"

# Function to read text content from an Excel document from Azure Blob Storage
def read_excel_document(document_path):
    try:
        blob_client = container_client.get_blob_client(document_path)
        with io.BytesIO(blob_client.download_blob().readall()) as file:
            workbook = openpyxl.load_workbook(io.BytesIO(file.read()))
            sheet = workbook.active
            content = ""
            for row in sheet.iter_rows(values_only=True):
                content += " ".join(map(str, row)) + "\n"
            return content
    except Exception as ex:
        return f"Error reading Excel document: {str(ex)}"

# Function to read text content from an ODS document from Azure Blob Storage
def read_ods_document(document_path):
    try:
        blob_client = container_client.get_blob_client(document_path)
        with io.BytesIO(blob_client.download_blob().readall()) as file:
            doc = ezodf.opendoc(io.BytesIO(file.read()))
            full_text = []
            for sheet in doc.sheets:
                for row in sheet.rows():
                    for cell in row:
                        if cell.value is not None:
                            if isinstance(cell.value, float):
                                full_text.append(str(cell.value))
                            else:
                                full_text.append(cell.value)
            return '\n'.join(full_text)
    except Exception as ex:
        return f"Error reading ODS document: {str(ex)}"

# Function to read text content from an ODT document from Azure Blob Storage
def read_odt_document(document_path):
    try:
        blob_client = container_client.get_blob_client(document_path)
        with io.BytesIO(blob_client.download_blob().readall()) as file:
            doc = ezodf.opendoc(io.BytesIO(file.read()))
            full_text = []
            for body_element in doc.body:
                if body_element.__class__.__name__ == "Text":
                    full_text.append(body_element.getvalue())
            return '\n'.join(full_text)
    except Exception as ex:
        return f"Error reading ODT document: {str(ex)}"

# Function to read text content from an ODP document from Azure Blob Storage
def read_odp_document(document_path):
    try:
        blob_client = container_client.get_blob_client(document_path)
        with io.BytesIO(blob_client.download_blob().readall()) as file:
            presentation = Presentation(io.BytesIO(file.read()))
            full_text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        full_text.append(shape.text)
            return '\n'.join(full_text)
    except Exception as ex:
        return f"Error reading ODP document: {str(ex)}"

# Function to combine and convert text from all documents into .txt format and store in Azure Blob Storage
def combine_and_convert_to_txt(documents, output_txt):
    combined_text = ""
    for document_path, document_content in documents.items():
        combined_text += f"Document: {document_path}\nContent: {document_content}\n\n"

    # Write combined text to Azure Blob Storage
    blob_client = container_client.get_blob_client(output_txt)
    blob_client.upload_blob(combined_text, overwrite=True)

    print(f"Combined text written to {output_txt} in Azure Blob Storage.")

    # Add the combined text to the documents dictionary
    documents[output_txt] = combined_text

# Function to get a description of the contents of a folder in Azure Blob Storage
def get_folder_contents_description(folder_path):
    try:
        blob_list = container_client.list_blobs(name_starts_with=folder_path)
        file_list = [blob.name for blob in blob_list]
        return f"The folder '{folder_path}' contains {len(file_list)} files: {', '.join(file_list)}"
    except Exception as ex:
        return f"Error getting folder contents description: {str(ex)}"

# Use the fine-tuned BERT model to send the query to OpenAI for analysis and response
def send_to_openai_and_get_response(input_zip, user_query, document_content, is_generic=False, document_path=None, trained_model_path=None):
    try:
        def classify_query(query):
            inputs = tokenizer(query, return_tensors="pt", padding=True, truncation=True)
            with torch.no_grad():
                outputs = model(**inputs)
            logits = outputs.logits
            probabilities = torch.softmax(logits, dim=1)
            relevant_probability = probabilities[0][1].item()
            return relevant_probability > 0.5  # Assuming a threshold of 0.5 for relevance classification

        if is_generic:
            if "type of files" in user_query.lower():
                return get_folder_contents_description(os.path.dirname(document_path))
            elif "document contain" in user_query.lower():
                return document_content
            elif "brief description" in user_query.lower():
                return get_folder_contents_description(os.path.dirname(document_path))
            elif "display contents" in user_query.lower():
                return get_folder_contents_description(os.path.dirname(document_path))
            elif "Hi" or 'hi' in user_query.lower():
                return "Hello ! How can I assist you today ?"
            else:
                return "Please provide a specific query related to the document content."
        else:
            if classify_query(user_query):
                response = openai.ChatCompletion.create(
                    model="gpt-3.5-turbo",
                    messages=[
                        {"role": "user", "content": user_query},
                        {"role": "assistant", "content": document_content}
                    ]
                )
                return response.choices[0].message['content'].strip()
            elif not user_query.strip():
                 print("You did not enter a query. Null value not accepted.")
            else:
                return "Query is not relevant to document content. No data available"
    except Exception as ex:
        return f"Error interacting with OpenAI: {str(ex)}"

output_zip = "combined_text.zip"  # Define output_zip before using it in the route function

@app.route('/bert_query', methods=['POST'])
def bert_query():
    try:
        data = request.get_json()  # Retrieve JSON data from the request
        
        # Access 'user_query' directly
        user_query = data.get('user_query')

        if not user_query:
            raise ValueError("Missing 'user_query' field in the form data.")

        # Your existing logic to process user input
        # Determine if the query is generic or document-specific
        is_generic = False
        user_queries = user_query.split(",")
        if len(user_queries) == 1 and user_queries[0].strip().lower().startswith("analyze document"):
            is_generic = True
        
        # Assuming you have the required variables (output_zip, output_txt, documents, folder_path, trained_model_path) defined elsewhere
        response = send_to_openai_and_get_response(output_zip, user_query, documents[output_txt],
                                                   is_generic=is_generic, document_path=folder_path,
                                                   trained_model_path=trained_model_path)
        print(response)
        return jsonify({"response": response})

    except Exception as e:
        # Log the error for debugging
        print(f"Error in process_query: {str(e)}")
        return jsonify({"error": str(e)})

if __name__ == "__main__":
    folder_path = "UploadedDocuments"

    documents = {}
    for blob in container_client.list_blobs():
        blob_name = blob.name
        if blob_name.endswith('.docx'):
            documents[blob_name] = read_word_document(blob_name)
        elif blob_name.endswith('.xlsx'):
            documents[blob_name] = read_excel_document(blob_name)
        elif blob_name.endswith('.pdf'):
            documents[blob_name] = read_pdf_document(blob_name)
        elif blob_name.endswith('.pptx'):
            documents[blob_name] = read_powerpoint_document(blob_name)
        elif blob_name.endswith('.ods'):
            documents[blob_name] = read_ods_document(blob_name)
        elif blob_name.endswith('.odt'):
            documents[blob_name] = read_odt_document(blob_name)
        elif blob_name.endswith('.odp'):
            documents[blob_name] = read_odp_document(blob_name)
        else:
            documents[blob_name] = read_text_document(blob_name)

    # Combine and convert text from all documents into .txt format and store in Azure Blob Storage
    output_txt = "combined_text.txt"
    combine_and_convert_to_txt(documents, output_txt)

    # Set your fine-tuned BERT model path
    fine_tuned_bert_model_path = "fine_tuned_document_search_model"

    # Path to the trained model CSV file
    trained_model_path = "training_set.csv"
    
    app.run(debug=True, port=5002)
